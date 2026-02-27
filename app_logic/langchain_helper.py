from langchain_google_genai import ChatGoogleGenerativeAI
from dotenv import load_dotenv

from langchain_core.prompts import PromptTemplate # VERY useful for structuring prompts
from langchain_core.output_parsers import StrOutputParser # VERY useful for parsing responses
from langchain_ollama.llms import OllamaLLM

# from langchain_community.agent_toolkits.load_tools import load_tools 
# from langchain.agents import create_agent
# from langchain_core.messages import HumanMessage, SystemMessage, AIMessage

import pypdf
import docx2txt
import io
import pytesseract
from PIL import Image
from pptx import Presentation

from langchain_text_splitters import RecursiveCharacterTextSplitter
import vector_helper as vh
import logging

load_dotenv()

logging.getLogger("pypdf").setLevel(logging.ERROR)

# TODO: Make an OCR to read files with images and images itself to convert to text
def ocr_image_to_text(image_file):
    try:
        img = Image.open(image_file)
        text = pytesseract.image_to_string(img)
        return text
    except Exception as e:
        print(f"OCR Error: {e}")
        return ""

def get_file_text(uploaded_file):
    if uploaded_file is None:
        return ""
    
    file_extension = uploaded_file.name.split('.')[-1].lower()

    if file_extension == 'pdf':
        reader = pypdf.PdfReader(uploaded_file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n" or ""
        return text
    
    elif file_extension == 'docx':
        return docx2txt.process(uploaded_file)
    
    elif file_extension == 'txt':
        return str(uploaded_file.read(), "utf-8")
    
    elif file_extension == 'md':
        stream = io.StringIO(uploaded_file.getvalue().decode("utf-8"))
        return stream.read()
    
    elif file_extension == 'pptx':
        prs = Presentation(uploaded_file)

        text_runs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)    

        return "\n".join(text_runs)
    
    elif file_extension in ['png', 'jpg', 'jpeg']:
        return ocr_image_to_text(uploaded_file)

    else:
        return ""
    
# Chunking text 
def chunk_text(text, chunk_size=1024, chunk_overlap=200):
    """
    :param text: Description
    :param chunk_size: Description
    :param chunk_overlap: Description
    :return: chunks
    :rtype: Any
    """

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", " ", ""],
        length_function=len,
    )

    chunks = text_splitter.split_text(text)
    return chunks

    # SIMULATION (including get_file_text method)
    """
    INPUTS:
    file = uploaded_file (Streamlit file uploader object)
    text = get_file_text(file)
    chunks = chunk_text(text, chunk_size=1024, chunk_overlap=200)

    OUTPUTS:
    chunks = [
        "A vegetable is a plant or part of a plant that is used as food. Examples include carrots, spinach, and potatoes.",
        "Fruits are typically sweet and contain seeds, while vegetables are more savory and do not contain seeds.",
        "Vegetables can be eaten raw or cooked, and they are an important source of vitamins and minerals in the diet.",
        "Common cooking methods for vegetables include steaming, boiling, roasting, and sautéing."
    ]
    # using vector store retrieval on chunks for RAG


    """
# TODO: make built in input templates for 7 diff quiz types and topics (DSA (lec, code), Stats (lec), OOP (lec, code), ATFL (lec), MS (lec) ) 
templates = [
    [

    ]
]


def gen_quiz(role, topics, context, raw_text, test_type, num_questions, difficulty, format, request, k):

    llm = ChatGoogleGenerativeAI(
        model="gemini-2.5-flash", 
        temperature=0.7
        )
    # llm = OllamaLLM(model="llama3.2", temperature=0.7)

    if raw_text:
        smart_context = vh.get_smart_context(raw_text, topics, k)
        final_context = f"User Context: {context}\n\nRelevant Context (from uploaded files): {smart_context}"
    else:
        final_context = context

    template_text = """
You are a {role}.

Retrieved text context: 
{context}

Topics: {topics}

Test types: {test_type} 
No. of items: {num_questions} 
Difficulty level: {difficulty}
Desired output or format: {format}

Generate a quiz based on the above information. {request}

DO NOT ADD ANY INTROS or any unnecessary things (just go straight to generating the quiz, that is all)
ONLY generate the quiz based on the information provided above.
IF YOU DO NOT HAVE ENOUGH INFORMATION, JUST SAY: I don't have relevant information about the {topics}, so I can't generate a quiz, PLEASE TRY RELEVANT TOPICS.
"""
    prompt_template_quiz = PromptTemplate(
        input_variables=[
                         'role', 
                         'topics', 
                         'context',  
                         'test_type', 
                         'num_questions', 
                         'difficulty', 
                         'format', 
                         'request'],

        template=template_text
    )

    formatted_prompt = prompt_template_quiz.format(
        role=role,
        topics=topics,
        context=final_context,
        test_type=test_type,
        num_questions=num_questions,
        difficulty=difficulty,
        format=format,
        request=request
    )

    chain = prompt_template_quiz | llm
    # chain = prompt_template_quiz | llm | StrOutputParser()  

    response = chain.invoke({
        "role": role,
        "topics": topics,
        "context": final_context,
        "test_type": test_type,
        "num_questions": num_questions,
        "difficulty": difficulty,
        "format": format,
        "request": request,
        "k": k
    })

    return response.content, formatted_prompt
    # return response, formatted_prompt


# testing the OLlama  (to be updated later)
def gen_quiz_local(role, topics, context, raw_text, test_type, num_questions, difficulty, format, request, k):
    llm = OllamaLLM(model="llama3.2", temperature=0.7)

    if raw_text:
        smart_context = vh.get_smart_context(raw_text, topics, k)
        final_context = f"User Context: {context}\n\nRelevant Context (from uploaded files): {smart_context}"
    else:
        final_context = context

    template_text = """
You are a {role}.

Retrieved text context: 
{context}

Topics: {topics}

Test types: {test_type} 
No. of items: {num_questions} 
Difficulty level: {difficulty}
Desired output or format: {format}

Generate a quiz based on the above information. {request}

DO NOT ADD ANY INTROS or any unnecessary things (just go straight to generating the quiz, that is all)
ONLY generate the quiz based on the information provided above.
IF YOU DO NOT HAVE ENOUGH INFORMATION, JUST SAY: I don't have relevant information about the {topics}, so I can't generate a quiz, PLEASE TRY RELEVANT TOPICS.
"""

    prompt_template_quiz = PromptTemplate(
        input_variables=['role', 
                         'topics', 
                         'context',
                         'extracted_text',
                         'test_type', 
                         'num_questions', 
                         'difficulty', 
                         'format', 
                         'request'],

        template= template_text
    )

    formatted_prompt = prompt_template_quiz.format(
        role=role,
        topics=topics,
        context=final_context,
        test_type=test_type,
        num_questions=num_questions,
        difficulty=difficulty,
        format=format,
        request=request
    )


    chain = prompt_template_quiz | llm | StrOutputParser()  

    response = chain.invoke({
        "role": role,
        "topics": topics,
        "context": final_context,
        "test_type": test_type,
        "num_questions": num_questions,
        "difficulty": difficulty,
        "format": format,
        "request": request
    })

    return response, formatted_prompt





# Testing
if __name__ == "__main__":
    # print("Pet Name Generation Example:")
    # print(generate_pet_name("cat", "orange", "playful"))
    print("\nQuiz Generation Example:")
    print(gen_quiz(
        role="history teacher",
        topics="World War II, The Renaissance",
        context="This quiz is for high school students studying European history.",
        raw_text=None,
        test_type="multiple choice and short answer",
        num_questions=5,
        difficulty="medium",
        format="JSON format with questions and answers",
        request="Please provide the correct answers as well.",
        k=5
    ))

    # print("\nLocal LLM Quiz Generation Example:")
    # print(gen_quiz_local(
    #     role="history teacher",
    #     topics="World War II, The Renaissance",
    #     context="This quiz is for high school students studying European history.",
    #     file=None,
    #     test_type="multiple choice and short answer",
    #     num_questions=5,
    #     difficulty="medium",
    #     format="JSON format with questions and answers",
    #     request="Please provide the correct answers as well."
    # ))

    
    # print("\nLangChain Agent Example:")
    # try:
    #     print(langchain_agent())
    # except Exception as e:
    #     print(f"Error: {e}")

# oh my god was I real 
# or someone who wore a plastic 
# smile and faded when no one 
# looked back at me again


"""
Summary of my RAG Pipeline:

1. Ingest: Files  --> String.
2. Chunk: String  --> Chunks.
3. Embed: Chunks  --> Vectors.
4. Retrieve: Topic  --> Search Vectors  --> Filter by Score.
5. Generate: Prompt + Smart Context  --> Quiz.
"""



# def langchain_agent():
#     llm = ChatGoogleGenerativeAI(
#         model="models/gemini-2.5-flash",
#         temperature=0.5
#     )

#     tools = load_tools(["wikipedia", "llm-math"], llm=llm)

#     system_prompt = SystemMessage(
#         content="You are a helpful and expert assistant. Use the provided tools to answer the question, then state the final answer clearly."
#     )

#     agent = create_agent(
#         llm, 
#         tools,
#         system_prompt=system_prompt 
#     )

#     inputs = {
#             "messages": [HumanMessage(
#                 content="What is the population of Canada and what is it divided by 3?"
#             )]
#         }

#     result = agent.invoke(inputs)

#     final_message_object = result['messages'][-1]

#     if isinstance(final_message_object, AIMessage):
#         final_message = final_message_object.content
#     else:
#         final_message = str(final_message_object)

#     return final_message